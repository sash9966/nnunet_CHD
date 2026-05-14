"""
Generate CHD Topology Loss presentation as PowerPoint.
Run: python3 scripts/make_presentation.py
Output: docs/CHD_TopologyLoss_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from pptx.oxml.ns import qn
from lxml import etree
import copy, math, os

# ─── COLOUR PALETTE ──────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x11, 0x17)   # slide background
BG_CARD   = RGBColor(0x16, 0x1B, 0x27)   # card / panel
BG_CODE   = RGBColor(0x1C, 0x21, 0x33)   # code block
BORDER    = RGBColor(0x2C, 0x31, 0x52)
PURPLE    = RGBColor(0x7C, 0x6A, 0xF7)   # primary accent (Claude purple)
CYAN      = RGBColor(0x4D, 0xC9, 0xE6)   # secondary accent
GREEN     = RGBColor(0x22, 0xD8, 0x7A)
YELLOW    = RGBColor(0xF5, 0xC5, 0x42)
RED       = RGBColor(0xF5, 0x65, 0x65)
ORANGE    = RGBColor(0xF6, 0xA6, 0x23)
WHITE     = RGBColor(0xE2, 0xE8, 0xF0)
MUTED     = RGBColor(0x88, 0x92, 0xB0)
DIM       = RGBColor(0x3C, 0x44, 0x62)

# Slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def rgb_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = font_name
    return txb

def add_text_box(slide, text, x, y, w, h,
                 size=Pt(13), bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, fill=None, line=None,
                 pad=Inches(0.15), font_name="Calibri"):
    """Text inside a rounded-rect card."""
    shape = add_rect(slide, x, y, w, h, fill=fill, line=line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = pad
    tf.margin_right = pad
    tf.margin_top   = pad
    tf.margin_bottom = pad
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return shape

def add_multiline(slide, lines, x, y, w, h,
                  default_size=Pt(13), fill=None, line=None, pad=Inches(0.18)):
    """lines: list of (text, size, bold, color, align)"""
    if fill or line:
        shape = add_rect(slide, x, y, w, h, fill=fill, line=line)
        tf = shape.text_frame
    else:
        txb = slide.shapes.add_textbox(x, y, w, h)
        tf  = txb.text_frame
    tf.word_wrap = True
    if fill or line:
        tf.margin_left = tf.margin_right = pad
        tf.margin_top  = tf.margin_bottom = pad

    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color, align = item, default_size, False, WHITE, PP_ALIGN.LEFT
        else:
            text = item[0]
            size  = item[1] if len(item) > 1 else default_size
            bold  = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else WHITE
            align = item[4] if len(item) > 4 else PP_ALIGN.LEFT

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = size
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = "Calibri"
    return tf

def add_bullet(slide, items, x, y, w, h,
               title=None, title_color=CYAN,
               size=Pt(12.5), fill=BG_CARD, line=BORDER,
               bullet_color=PURPLE, text_color=WHITE):
    shape = add_rect(slide, x, y, w, h, fill=fill, line=line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top  = tf.margin_bottom = Inches(0.15)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(13)
        run.font.bold  = True
        run.font.color.rgb = title_color
        run.font.name  = "Calibri"
        first = False
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        run = p.add_run()
        run.text = "◆  "
        run.font.size  = Pt(9)
        run.font.color.rgb = bullet_color
        run.font.bold  = True
        run.font.name  = "Calibri"
        run2 = p.add_run()
        run2.text = item
        run2.font.size  = size
        run2.font.color.rgb = text_color
        run2.font.name  = "Calibri"
    return shape

def slide_header(slide, title, subtitle=None, accent=PURPLE):
    # Top colour bar
    add_rect(slide, 0, 0, W, Inches(0.06), fill=accent)
    # Title
    add_text(slide, title,
             Inches(0.55), Inches(0.14), W - Inches(1.1), Inches(0.55),
             size=Pt(22), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.55), Inches(0.64), W - Inches(1.1), Inches(0.36),
                 size=Pt(13), color=MUTED, italic=True)
    # Divider line
    add_rect(slide, Inches(0.55), Inches(0.98), W - Inches(1.1), Inches(0.02), fill=BORDER)

def slide_number(slide, n, total=14):
    add_text(slide, f"{n} / {total}",
             W - Inches(1.0), H - Inches(0.35), Inches(0.85), Inches(0.28),
             size=Pt(9), color=DIM, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
def make_title_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)

    # Large gradient-like decoration rectangles
    add_rect(sl, W - Inches(4.5), 0, Inches(4.5), H,
             fill=RGBColor(0x10, 0x14, 0x22))
    add_rect(sl, W - Inches(4.6), Inches(1.6), Inches(0.08), Inches(4.3),
             fill=PURPLE)
    # Purple glow block top-right
    add_rect(sl, W - Inches(3.8), Inches(0.3), Inches(3.5), Inches(2.0),
             fill=RGBColor(0x1A, 0x16, 0x38))
    add_rect(sl, W - Inches(3.8), Inches(0.3), Inches(3.5), Inches(0.04),
             fill=PURPLE)

    # Title text
    add_text(sl, "Topology-Preserving Segmentation",
             Inches(0.6), Inches(1.5), Inches(7.8), Inches(0.8),
             size=Pt(34), bold=True, color=WHITE)
    add_text(sl, "of Congenital Heart Disease",
             Inches(0.6), Inches(2.2), Inches(7.8), Inches(0.7),
             size=Pt(34), bold=True, color=PURPLE)

    add_text(sl, "Cascaded nnU-Net with FiLM Disease Conditioning & soft-clDice Loss",
             Inches(0.6), Inches(3.05), Inches(7.5), Inches(0.5),
             size=Pt(15), color=CYAN, italic=True)

    # Divider
    add_rect(sl, Inches(0.6), Inches(3.65), Inches(4.5), Inches(0.03), fill=BORDER)

    add_text(sl, "nnU-Net V2 fork  ·  all-experiments branch  ·  Stanford / Akshay Lab",
             Inches(0.6), Inches(3.78), Inches(7.5), Inches(0.35),
             size=Pt(11), color=MUTED)

    # Right panel: quick facts
    facts = [
        ("22", "Composed Trainers"),
        ("5",  "Mixin Features"),
        ("8",  "Cardiac Labels"),
        ("8",  "Disease Flags"),
    ]
    for i, (num, lbl) in enumerate(facts):
        cx = W - Inches(3.5) + Inches(i % 2) * Inches(1.75)
        cy = Inches(2.5) + (i // 2) * Inches(1.5)
        add_rect(sl, cx, cy, Inches(1.6), Inches(1.25),
                 fill=RGBColor(0x1C, 0x20, 0x38), line=BORDER)
        add_text(sl, num, cx, cy + Inches(0.12), Inches(1.6), Inches(0.65),
                 size=Pt(32), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, cx, cy + Inches(0.7), Inches(1.6), Inches(0.4),
                 size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

    slide_number(sl, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — MOTIVATION
# ─────────────────────────────────────────────────────────────────────────────
def make_motivation_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Motivation",
                 "Why do we need topology loss and disease conditioning?")
    slide_number(sl, 2)

    # Baseline stat
    add_rect(sl, Inches(0.5), Inches(1.2), Inches(2.2), Inches(1.0),
             fill=RGBColor(0x1A, 0x22, 0x16), line=GREEN)
    add_text(sl, "~0.95", Inches(0.5), Inches(1.25), Inches(2.2), Inches(0.5),
             size=Pt(28), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "Whole-Heart Dice\n(3d_fullres baseline)",
             Inches(0.5), Inches(1.68), Inches(2.2), Inches(0.45),
             size=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER)

    add_text(sl, "→  But two systematic failure modes remain:",
             Inches(2.9), Inches(1.38), Inches(6), Inches(0.4),
             size=Pt(13), color=CYAN, bold=True)

    # Failure 1
    add_bullet(sl,
        ["AO/PA label confusion in TGA & DORV cases",
         "Fullres patch lacks global context to determine which vessel exits which ventricle",
         "TGA: aorta exits RV, PA exits LV — reversed from normal anatomy",
         "DORV: both great vessels exit RV — both labels spatially ambiguous"],
        Inches(0.5), Inches(2.0), Inches(5.9), Inches(1.8),
        title="Failure 1 — Great Vessel Label Confusion",
        title_color=RED, bullet_color=RED)

    # Failure 2
    add_bullet(sl,
        ["VSD: inter-ventricular septum is absent or incomplete",
         "Model predicts single large LV region instead of LV + RV",
         "AVSD: all four-chamber septal borders ambiguous simultaneously"],
        Inches(0.5), Inches(3.95), Inches(5.9), Inches(1.6),
        title="Failure 2 — VSD Boundary Ambiguity",
        title_color=YELLOW, bullet_color=YELLOW)

    # Solutions panel
    add_bullet(sl,
        ["Cascade pipeline: low-res stage sees full cardiac anatomy in one patch → global context",
         "FiLM conditioning: inject disease vector to bias network toward correct anatomy",
         "soft-clDice topology loss: penalise broken vessel centerlines on AO & PA"],
        Inches(7.0), Inches(1.2), Inches(5.8), Inches(4.5),
        title="Our Solutions",
        title_color=PURPLE, bullet_color=PURPLE)

make_title_slide()
make_motivation_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — LABEL & DISEASE MAPPING
# ─────────────────────────────────────────────────────────────────────────────
def make_labels_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Data: Labels & Disease Flags",
                 "8 cardiac structures  ·  K=8 binary disease flags per case")
    slide_number(sl, 3)

    # LEFT: label table
    Y0 = Inches(1.18)
    headers = ["ID", "Structure", "Abbrev", "Topo?"]
    col_w   = [Inches(0.5), Inches(2.5), Inches(0.85), Inches(0.75)]
    col_x   = [Inches(0.4)]
    for w in col_w[:-1]: col_x.append(col_x[-1] + w)
    ROW_H = Inches(0.38)

    # Header row
    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_rect(sl, cx, Y0, cw, ROW_H, fill=RGBColor(0x20, 0x22, 0x40), line=BORDER)
        add_text(sl, hdr, cx, Y0 + Inches(0.07), cw, ROW_H,
                 size=Pt(10), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    rows = [
        ("0","Background","BG",None),
        ("1","LV Blood Pool","LV-BP",None),
        ("2","RV Blood Pool","RV-BP",None),
        ("3","Left Atrium","LA",None),
        ("4","Right Atrium","RA",None),
        ("5","Myocardium","Myo",None),
        ("6","Aorta","AO","clDice ✓"),
        ("7","Pulmonary A.","PA","clDice ✓"),
    ]
    for r, (rid, name, abbr, topo) in enumerate(rows):
        ry = Y0 + ROW_H * (r + 1)
        bg = RGBColor(0x1A, 0x1E, 0x32) if r % 2 == 0 else BG_DARK
        for i, (val, cx, cw) in enumerate(zip([rid, name, abbr, topo or "—"], col_x, col_w)):
            add_rect(sl, cx, ry, cw, ROW_H, fill=bg, line=BORDER)
            c = PURPLE if topo and i == 3 else (CYAN if i == 0 else WHITE)
            add_text(sl, val, cx + Inches(0.04), ry + Inches(0.08), cw - Inches(0.08), ROW_H,
                     size=Pt(10.5), color=c, align=PP_ALIGN.CENTER)

    # RIGHT: disease vector
    DX = Inches(5.05)
    disease = [
        ("0","HLHS","LV may be absent/tiny; AO unreliable"),
        ("1","ASD","LA–RA may be adjacent"),
        ("2","VSD","LV–RV boundary ambiguous"),
        ("3","AVSD","All 4-chamber borders ambiguous"),
        ("4","DORV","Both AO & PA exit RV"),
        ("5","PuA","AO/PA fused — vessel absent"),
        ("6","ToF","Overriding aorta straddles VSD"),
        ("7","TGA","AO↔RV, PA↔LV (reversed)"),
    ]
    add_text(sl, "disease_map.json  →  {\"SV_001\": [0,0,1,0,0,0,0,1], ...}",
             DX, Inches(1.1), Inches(7.8), Inches(0.35),
             size=Pt(10), color=RGBColor(0xA8,0xD8,0xFF), font_name="Courier New")

    DCX = [DX, DX + Inches(0.55), DX + Inches(1.3), DX + Inches(5.3)]
    DCW = [Inches(0.5), Inches(0.7), Inches(3.9), Inches(2.25)]
    dhdrs = ["Idx","Flag","Key implication","Topology risk"]
    drisk = [MUTED,MUTED,MUTED,MUTED,MUTED,MUTED,ORANGE,RED]
    drisk_txt = ["Low","Low","VSD","AVSD","DORV","Low","ToF","TGA ★"]

    for i, (hdr, cx, cw) in enumerate(zip(dhdrs, DCX, DCW)):
        add_rect(sl, cx, Y0, cw, ROW_H, fill=RGBColor(0x20,0x22,0x40), line=BORDER)
        add_text(sl, hdr, cx, Y0 + Inches(0.07), cw, ROW_H,
                 size=Pt(10), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    for r, (did, flag, impl) in enumerate(disease):
        ry = Y0 + ROW_H * (r + 1)
        bg = RGBColor(0x1A,0x1E,0x32) if r % 2 == 0 else BG_DARK
        for i, (val, cx, cw) in enumerate(zip([did, flag, impl, drisk_txt[r]], DCX, DCW)):
            add_rect(sl, cx, ry, cw, ROW_H, fill=bg, line=BORDER)
            c = CYAN if i == 1 else (drisk[r] if i == 3 else WHITE)
            add_text(sl, val, cx + Inches(0.04), ry + Inches(0.07), cw - Inches(0.08), ROW_H,
                     size=Pt(10), color=c, align=PP_ALIGN.CENTER if i < 2 else PP_ALIGN.LEFT)

make_labels_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — CASCADE PIPELINE
# ─────────────────────────────────────────────────────────────────────────────
def make_cascade_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Cascade Pipeline Architecture",
                 "Two-stage approach: global context → local boundary refinement")
    slide_number(sl, 4)

    # Stage 1 box
    def stage_box(x, y, w, h, title, color, items):
        add_rect(sl, x, y, w, h, fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_rect(sl, x, y, w, Inches(0.05), fill=color)  # top bar
        add_text(sl, title, x + Inches(0.12), y + Inches(0.1), w - Inches(0.24), Inches(0.4),
                 size=Pt(13), bold=True, color=color)
        for i, item in enumerate(items):
            add_text(sl, f"◆  {item}",
                     x + Inches(0.12), y + Inches(0.52 + i * 0.38),
                     w - Inches(0.24), Inches(0.36),
                     size=Pt(11), color=WHITE)

    stage_box(Inches(0.4), Inches(1.2), Inches(4.0), Inches(4.2),
              "Stage 1 — 3d_lowres", CYAN,
              ["~3–4 mm voxel spacing",
               "Full heart in every patch",
               "Global label assignment",
               "FiLM: disease-aware topology",
               "Optional: clDice on AO/PA",
               "Output → coarse labels"])

    # Arrow
    add_text(sl, "→", Inches(4.5), Inches(2.9), Inches(0.5), Inches(0.5),
             size=Pt(36), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    # Concatenation block
    add_rect(sl, Inches(5.1), Inches(2.7), Inches(1.2), Inches(0.9),
             fill=RGBColor(0x1E, 0x18, 0x38), line=PURPLE)
    add_text(sl, "concat\none-hot\nprior",
             Inches(5.1), Inches(2.72), Inches(1.2), Inches(0.85),
             size=Pt(9), color=PURPLE, align=PP_ALIGN.CENTER)

    # Arrow 2
    add_text(sl, "→", Inches(6.4), Inches(2.9), Inches(0.5), Inches(0.5),
             size=Pt(36), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    stage_box(Inches(7.0), Inches(1.2), Inches(4.2), Inches(4.2),
              "Stage 2 — 3d_cascade_fullres", PURPLE,
              ["~1 mm voxel spacing",
               "Local boundary crop",
               "Prior from Stage 1 as extra channels",
               "FiLM: disease-aware refinement",
               "Optional: clDice on AO/PA",
               "Output → final labels"])

    # Bottom note
    add_text(sl, "Key insight:  Stage 1 sees the entire cardiac anatomy → correct global vessel assignment.  "
                 "Stage 2 sharpens boundaries using the Stage 1 prior + full resolution.",
             Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.6),
             size=Pt(11.5), color=MUTED, italic=True)

    # Note on planner
    add_text(sl,
             "⚠  ExperimentPlannerForceLowRes required — default nnU-Net planner skips 3d_lowres for cardiac CT",
             Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.4),
             size=Pt(10.5), color=YELLOW)

make_cascade_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — FILM CONDITIONING
# ─────────────────────────────────────────────────────────────────────────────
def make_film_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "FiLM Disease Conditioning",
                 "Feature-wise Linear Modulation at the bottleneck")
    slide_number(sl, 5)

    # Architecture diagram (text-based)
    arch = [
        ("Image → Encoder → Bottleneck", Pt(12), False, MUTED),
        ("", Pt(6), False, MUTED),
        ("          ↓", Pt(12), False, WHITE),
        ("          disease_mlp( vec_K8 )  →  γ, β  (channel-wise)", Pt(11.5), False, CYAN),
        ("          feature  =  (1 + γ) × feature + β", Pt(11.5), True, WHITE),
        ("          ↓", Pt(12), False, WHITE),
        ("Decoder → Segmentation", Pt(12), False, MUTED),
    ]
    add_multiline(sl, [(t, s, b, c) for t, s, b, c in arch],
                  Inches(0.4), Inches(1.2), Inches(5.4), Inches(2.8),
                  fill=BG_CODE, line=BORDER)

    # Key design decisions
    decisions = [
        ("Bottleneck-only", PURPLE,
         "Decoder has N=7 stages → (1+γ)^7 ≈ 2× feature distortion with γ=0.1. Applied at bottleneck only."),
        ("Near-zero init", CYAN,
         "weights σ=0.01 → identity transform at start. Training begins from unconditioned baseline."),
        ("LR × 2.0", GREEN,
         "disease_mlp + bottleneck_film params use 2× learning rate to learn conditioning faster."),
        ("CFG dropout 10%", YELLOW,
         "10% of batches use zero disease vector. Prevents over-reliance on disease label at inference."),
    ]
    for i, (title, color, desc) in enumerate(decisions):
        x = Inches(0.4) + (i % 2) * Inches(3.05)
        y = Inches(4.15) + (i // 2) * Inches(1.35)
        add_rect(sl, x, y, Inches(2.85), Inches(1.2),
                 fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_rect(sl, x, y, Inches(2.85), Inches(0.04), fill=color)
        add_text(sl, title, x + Inches(0.12), y + Inches(0.08), Inches(2.6), Inches(0.35),
                 size=Pt(12), bold=True, color=color)
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.42), Inches(2.6), Inches(0.7),
                 size=Pt(10), color=MUTED)

    # Instability lesson
    add_rect(sl, Inches(6.1), Inches(1.2), Inches(6.7), Inches(5.2),
             fill=RGBColor(0x1A,0x14,0x2A), line=PURPLE)
    add_rect(sl, Inches(6.1), Inches(1.2), Inches(6.7), Inches(0.04), fill=PURPLE)
    add_text(sl, "Instability Lesson (Early Experiments)",
             Inches(6.25), Inches(1.28), Inches(6.4), Inches(0.4),
             size=Pt(13), bold=True, color=PURPLE)

    add_multiline(sl, [
        ("Applied FiLM at all 7 decoder stages →", Pt(11.5), False, WHITE),
        ("   (1 + γ)^7  ≈  1.95×  feature distortion  at  γ = 0.1", Pt(11), True, RED),
        ("", Pt(6), False, WHITE),
        ("With LR×10 + momentum 0.99:", Pt(11.5), False, WHITE),
        ("   γ grew large fast → training diverged", Pt(11), False, RED),
        ("", Pt(6), False, WHITE),
        ("Fix: bottleneck-only + LR×2 + near-zero init", Pt(11.5), True, GREEN),
        ("   Training is now stable across all experiments", Pt(11), False, MUTED),
        ("", Pt(6), False, WHITE),
        ("disease_param_prefixes = {", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("    'disease_mlp',", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("    'bottleneck_film'", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("}", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
    ], Inches(6.25), Inches(1.72), Inches(6.35), Inches(4.5), fill=None)

make_film_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — WHAT IS clDice
# ─────────────────────────────────────────────────────────────────────────────
def make_cldice_intro_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Topology Loss: soft-clDice",
                 "Shit et al., CVPR 2021 — clDice: A Novel Topology-Preserving Loss")
    slide_number(sl, 6)

    # Citation box
    add_rect(sl, Inches(0.4), Inches(1.15), Inches(7.5), Inches(0.55),
             fill=RGBColor(0x14,0x18,0x30), line=BORDER)
    add_text(sl, 'S. Shit et al.  "clDice - a Novel Topology-Preserving Loss Function for Tubular Structure Segmentation."  '
             'CVPR 2021.',
             Inches(0.55), Inches(1.2), Inches(7.2), Inches(0.45),
             size=Pt(10.5), color=MUTED, italic=True)

    # Core idea
    add_bullet(sl,
        ["Standard Dice measures voxel overlap — insensitive to broken connectivity",
         "Two predictions with same Dice can have very different topology (one vessel vs many fragments)",
         "clDice: compare soft skeletons (centerlines) of prediction and GT",
         "A broken vessel → skeleton fragments → low skeleton overlap → high loss",
         "Fully differentiable — gradients flow through soft morphological operations"],
        Inches(0.4), Inches(1.85), Inches(5.8), Inches(2.5),
        title="Core Idea", title_color=PURPLE, bullet_color=PURPLE)

    # Formula box
    add_multiline(sl, [
        ("clDice Formula", Pt(13), True, CYAN),
        ("", Pt(6), False, WHITE),
        ("tprec = Σ(skel_pred · gt)   / Σ(skel_pred)   ← skeleton inside GT", Pt(11), False, WHITE),
        ("tsens = Σ(skel_gt  · pred)  / Σ(skel_gt)     ← GT skeleton inside pred", Pt(11), False, WHITE),
        ("", Pt(6), False, WHITE),
        ("clDice = 2·tprec·tsens / (tprec + tsens + ε)", Pt(13), True, PURPLE),
        ("loss   = 1 − clDice", Pt(13), True, RED),
        ("", Pt(6), False, WHITE),
        ("Applied as extra term:  total = Dice+CE  +  w · clDice", Pt(11), False, YELLOW),
    ], Inches(0.4), Inches(4.5), Inches(5.8), Inches(2.55),
       fill=BG_CODE, line=BORDER)

    # Right: Why AO/PA only
    add_bullet(sl,
        ["AO and PA are tubular (vessel-like) — clDice designed for tubes",
         "Dominant failure modes (TGA, DORV) involve AO/PA topology",
         "Non-tubular structures (LV, RV, atria) do not benefit from skeleton loss",
         "Applying to all 8 classes × 10 skeleton iters would be computationally prohibitive",
         "Label IDs resolved dynamically from dataset.json keyword match"],
        Inches(6.5), Inches(1.85), Inches(6.3), Inches(2.5),
        title="Why Applied to AO (label 6) & PA (label 7) Only?",
        title_color=CYAN, bullet_color=CYAN)

    # 3D note
    add_rect(sl, Inches(6.5), Inches(4.5), Inches(6.3), Inches(2.55),
             fill=RGBColor(0x14, 0x22, 0x14), line=GREEN)
    add_rect(sl, Inches(6.5), Inches(4.5), Inches(6.3), Inches(0.04), fill=GREEN)
    add_text(sl, "3D Implementation — Fully Supported",
             Inches(6.65), Inches(4.58), Inches(6.0), Inches(0.38),
             size=Pt(12), bold=True, color=GREEN)
    add_text(sl,
             "soft_erode / soft_dilate detect tensor dimensionality:\n"
             "  4D  (B,1,H,W)   → F.max_pool2d  (kernel 3×3)\n"
             "  5D  (B,1,D,H,W) → F.max_pool3d  (kernel 3×3×3)\n\n"
             "No code changes needed for 3D — already dimension-agnostic.",
             Inches(6.65), Inches(5.0), Inches(6.0), Inches(1.9),
             size=Pt(11), color=WHITE)

make_cldice_intro_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — clDice ALGORITHM
# ─────────────────────────────────────────────────────────────────────────────
def make_algorithm_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Soft Skeletonization Algorithm",
                 "Differentiable skeleton via iterative morphological opening")
    slide_number(sl, 7)

    add_multiline(sl, [
        ("# Soft morphological primitives", Pt(11), False, MUTED),
        ("def soft_erode(x):   return -max_pool3d(-x, k=3, pad=1)   # min-pool", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("def soft_dilate(x):  return  max_pool3d( x, k=3, pad=1)   # max-pool", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("def soft_open(x):    return soft_dilate(soft_erode(x))     # open = erode→dilate", Pt(11), False, RGBColor(0xA8,0xD8,0xFF)),
        ("", Pt(7), False, WHITE),
        ("# Iterative skeleton accumulation", Pt(11), False, MUTED),
        ("skeleton = zeros_like(x)", Pt(11), False, WHITE),
        ("current  = x.clone()", Pt(11), False, WHITE),
        ("for _ in range(num_iter=10):", Pt(11), False, WHITE),
        ("    opened   = soft_open(current)           # remove thin structures", Pt(11), False, CYAN),
        ("    skeleton += ReLU(current - opened)      # residual = thin structures", Pt(11), True, PURPLE),
        ("    current   = soft_erode(current)         # shrink further", Pt(11), False, CYAN),
        ("return skeleton", Pt(11), False, WHITE),
    ], Inches(0.4), Inches(1.15), Inches(6.1), Inches(4.4), fill=BG_CODE, line=BORDER)

    # Right: properties
    add_bullet(sl,
        ["Differentiable — gradients backpropagate through all pooling ops",
         "GT skeleton computed with torch.no_grad() (no gradient needed)",
         "num_iter=10 controls skeleton thickness (higher → captures thicker tubes)",
         "Kernel size fixed at 3 (3×3×3 in 3D) — smallest meaningful structuring element",
         "Memory: 10 × 2 pooling ops per class per batch  →  2 classes (AO+PA) = 40 ops"],
        Inches(6.7), Inches(1.15), Inches(6.1), Inches(2.65),
        title="Algorithm Properties", title_color=PURPLE)

    # Iteration diagram (simplified visual)
    add_text(sl, "How the skeleton accumulates across iterations:",
             Inches(6.7), Inches(4.0), Inches(6.1), Inches(0.35),
             size=Pt(11.5), color=CYAN, bold=True)

    iters = [
        ("iter 0", "Full vessel mask", GREEN),
        ("iter 3", "Medial region only", YELLOW),
        ("iter 6", "Near centerline", ORANGE),
        ("iter 9", "Centerline voxels", RED),
        ("∑ residuals", "Soft skeleton", PURPLE),
    ]
    for i, (label, desc, color) in enumerate(iters):
        bx = Inches(6.7) + i * Inches(1.22)
        add_rect(sl, bx, Inches(4.4), Inches(1.05), Inches(0.85),
                 fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_text(sl, label, bx, Inches(4.42), Inches(1.05), Inches(0.32),
                 size=Pt(9.5), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(sl, desc, bx, Inches(4.72), Inches(1.05), Inches(0.5),
                 size=Pt(8.5), color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(iters) - 1:
            add_text(sl, "→", bx + Inches(1.05), Inches(4.62), Inches(0.17), Inches(0.35),
                     size=Pt(14), color=DIM, align=PP_ALIGN.CENTER)

    # skeletonize note
    add_text(sl,
             "Accumulating residuals across all iterations produces a soft approximation of the "
             "morphological skeleton — the set of centerline voxels equidistant from the vessel boundary.",
             Inches(6.7), Inches(5.45), Inches(6.1), Inches(0.7),
             size=Pt(10.5), color=MUTED, italic=True)

make_algorithm_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — WEIGHT SCHEDULE
# ─────────────────────────────────────────────────────────────────────────────
def make_schedule_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Two Topology Mixin Variants",
                 "Fixed weight (existing)  vs.  Scheduled weight (new)")
    slide_number(sl, 8)

    # Fixed mixin
    add_rect(sl, Inches(0.4), Inches(1.15), Inches(6.0), Inches(4.0),
             fill=RGBColor(0x14,0x18,0x2C), line=YELLOW)
    add_rect(sl, Inches(0.4), Inches(1.15), Inches(6.0), Inches(0.04), fill=YELLOW)
    add_text(sl, "TopologyLossMixin  (existing, all current trainers)",
             Inches(0.55), Inches(1.22), Inches(5.7), Inches(0.4),
             size=Pt(12.5), bold=True, color=YELLOW)

    add_multiline(sl, [
        ("topo_weight   =  0.2   (constant)", Pt(12), True, WHITE),
        ("topo_num_iter =  10", Pt(11.5), False, MUTED),
        ("topo_class_ids = [6, 7]   (AO, PA)", Pt(11.5), False, MUTED),
        ("", Pt(6), False, WHITE),
        ("Loss = Dice+CE + 0.2 × clDice", Pt(12), False, YELLOW),
        ("", Pt(6), False, WHITE),
        ("Applied from epoch 0 → topology loss active", Pt(11), False, WHITE),
        ("immediately, even when predictions are random.", Pt(11), False, WHITE),
        ("", Pt(6), False, WHITE),
        ("⚠  Risk: topology term may interfere with", Pt(11), False, YELLOW),
        ("    Dice+CE convergence in early epochs.", Pt(11), False, YELLOW),
    ], Inches(0.55), Inches(1.72), Inches(5.7), Inches(3.2), fill=None)

    # Scheduled mixin
    add_rect(sl, Inches(6.7), Inches(1.15), Inches(6.2), Inches(4.0),
             fill=RGBColor(0x14,0x18,0x2C), line=PURPLE)
    add_rect(sl, Inches(6.7), Inches(1.15), Inches(6.2), Inches(0.04), fill=PURPLE)
    add_text(sl, "TopologyLossScheduledMixin  (new)",
             Inches(6.85), Inches(1.22), Inches(5.9), Inches(0.4),
             size=Pt(12.5), bold=True, color=PURPLE)

    add_multiline(sl, [
        ("topo_warmup_epochs  =  10   (0 → w_high ramp)", Pt(11.5), False, WHITE),
        ("topo_decay_start    =  40   (plateau ends)", Pt(11.5), False, WHITE),
        ("topo_w_high         =  0.5  (peak weight)", Pt(11.5), True, PURPLE),
        ("topo_w_low          =  0.05 (final weight)", Pt(11.5), False, MUTED),
        ("", Pt(6), False, WHITE),
        ("Epoch  0–9 :  linear ramp  0.0 → 0.5", Pt(11), False, GREEN),
        ("Epoch 10–39:  plateau at   0.5", Pt(11), False, CYAN),
        ("Epoch 40–99:  cosine decay 0.5 → 0.05", Pt(11), False, YELLOW),
        ("", Pt(6), False, WHITE),
        ("✓  Dice+CE converges first (epochs 0-9)", Pt(11), False, GREEN),
        ("✓  Topology enforced at peak (epoch 10-39)", Pt(11), False, GREEN),
        ("✓  Gentle decay prevents over-constraint", Pt(11), False, GREEN),
    ], Inches(6.85), Inches(1.72), Inches(5.9), Inches(3.2), fill=None)

    # Weight chart (ASCII-style bar visual)
    add_text(sl, "Weight schedule over 100 epochs:",
             Inches(0.4), Inches(5.35), Inches(5.0), Inches(0.3),
             size=Pt(11), color=MUTED)

    phases = [
        ("Warmup\n(0→10)", Inches(0.4), Inches(0.75), GREEN),
        ("Plateau\n(10→40)", Inches(1.22), Inches(1.85), CYAN),
        ("Cosine Decay\n(40→100)", Inches(3.14), Inches(3.06), YELLOW),
    ]
    for label, bx, bw, color in phases:
        add_rect(sl, bx, Inches(5.7), bw, Inches(0.5),
                 fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_text(sl, label, bx, Inches(5.72), bw, Inches(0.46),
                 size=Pt(9), color=color, align=PP_ALIGN.CENTER)

    add_text(sl, "←─────────────────── 100 epochs ───────────────────→",
             Inches(0.4), Inches(6.28), Inches(6.0), Inches(0.3),
             size=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER)

    # Key difference note
    add_rect(sl, Inches(6.7), Inches(5.35), Inches(6.2), Inches(1.25),
             fill=RGBColor(0x16,0x12,0x2E), line=PURPLE)
    add_text(sl,
             "Both mixins share the same mixin_extra_loss() implementation "
             "(inherited from TopologyLossMixin). "
             "TopologyLossScheduledMixin only adds mixin_on_train_epoch_start() "
             "to update self.topo_weight each epoch using topo_weight_schedule().",
             Inches(6.85), Inches(5.42), Inches(5.9), Inches(1.1),
             size=Pt(10.5), color=MUTED, italic=True)

make_schedule_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — MIXIN ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
def make_mixin_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Composable Mixin Architecture",
                 "Feature modules stack via Python MRO — no monolithic subclassing")
    slide_number(sl, 9)

    # Mixin table
    mixins = [
        ("DiseaseConditioningMixin", CYAN,   "Injects disease vector via FiLM or MLP wrapper"),
        ("TopologyLossMixin",        PURPLE, "Adds soft-clDice on AO/PA  (fixed weight 0.2)"),
        ("TopologyLossScheduledMixin",PURPLE,"Like above but with warmup→plateau→decay schedule"),
        ("CurriculumWeightsMixin",   YELLOW, "Per-class CE weights that evolve during training"),
    ]
    add_text(sl, "Available Mixins", Inches(0.4), Inches(1.15), Inches(6.5), Inches(0.32),
             size=Pt(13), bold=True, color=WHITE)
    for i, (name, color, desc) in enumerate(mixins):
        y = Inches(1.52) + i * Inches(0.62)
        add_rect(sl, Inches(0.4), y, Inches(6.5), Inches(0.55),
                 fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_text(sl, name, Inches(0.55), y + Inches(0.06), Inches(3.2), Inches(0.25),
                 size=Pt(11.5), bold=True, color=color, font_name="Courier New")
        add_text(sl, desc, Inches(0.55), y + Inches(0.28), Inches(5.8), Inches(0.22),
                 size=Pt(10), color=MUTED)

    # Hooks table
    add_text(sl, "Hook Dispatch Chain  (each mixin calls super())",
             Inches(0.4), Inches(4.12), Inches(6.5), Inches(0.32),
             size=Pt(13), bold=True, color=WHITE)
    hooks = [
        ("mixin_init()",              "Set mixin-specific __init__ attributes"),
        ("mixin_initialize()",        "Post-setup: allocate GPU objects, resolve labels"),
        ("mixin_extra_loss(...)",      "Return extra loss scalar to add to main loss"),
        ("mixin_param_groups()",       "Extra optimizer param groups with custom LR"),
        ("mixin_on_train_epoch_start()","Per-epoch setup (update schedules, etc.)"),
    ]
    for i, (hook, desc) in enumerate(hooks):
        y = Inches(4.5) + i * Inches(0.42)
        add_rect(sl, Inches(0.4), y, Inches(2.65), Inches(0.38),
                 fill=BG_CODE, line=BORDER)
        add_text(sl, hook, Inches(0.5), y + Inches(0.06), Inches(2.5), Inches(0.26),
                 size=Pt(10), bold=True, color=CYAN, font_name="Courier New")
        add_text(sl, desc, Inches(3.15), y + Inches(0.08), Inches(3.7), Inches(0.26),
                 size=Pt(10), color=MUTED)

    # MRO examples
    add_text(sl, "Example MRO Chains", Inches(7.1), Inches(1.15), Inches(5.8), Inches(0.32),
             size=Pt(13), bold=True, color=WHITE)

    examples = [
        ("DA5CascadeFiLMTopo  (low-res)",
         ["DiseaseConditioningMixin", "TopologyLossMixin",
          "ComposableTrainerMixin", "nnUNetTrainerDA5", "TrainerMixin"],
         CYAN),
        ("DA5CascadeFullresFiLMTopo  (NEW fixed)",
         ["DiseaseConditioningMixin", "TopologyLossMixin",
          "ComposableTrainerMixin", "nnUNetTrainerDA5", "TrainerMixin"],
         PURPLE),
        ("DA5CascadeTopoScheduled  (NEW)",
         ["TopologyLossScheduledMixin", "TopologyLossMixin",
          "ComposableTrainerMixin", "nnUNetTrainerDA5", "TrainerMixin"],
         GREEN),
        ("DA5FiLMTopoCurriculum  (full-res)",
         ["DiseaseConditioningMixin", "TopologyLossMixin",
          "CurriculumWeightsMixin", "ComposableTrainerMixin", "nnUNetTrainerDA5"],
         YELLOW),
    ]
    for i, (label, chain, color) in enumerate(examples):
        y = Inches(1.55) + i * Inches(1.38)
        add_text(sl, label, Inches(7.1), y, Inches(5.8), Inches(0.3),
                 size=Pt(10.5), bold=True, color=color)
        # chain pills
        for j, mixin in enumerate(chain):
            mx = Inches(7.1) + j * Inches(1.18)
            add_rect(sl, mx, y + Inches(0.33), Inches(1.12), Inches(0.36),
                     fill=RGBColor(0x14,0x18,0x2C), line=color if j < 2 else BORDER)
            add_text(sl, mixin.replace("Mixin","").replace("nnUNetTrainer",""),
                     mx + Inches(0.03), y + Inches(0.35), Inches(1.06), Inches(0.3),
                     size=Pt(7.5), color=color if j < 2 else MUTED, align=PP_ALIGN.CENTER)
            if j < len(chain) - 1:
                add_text(sl, "→", mx + Inches(1.12), y + Inches(0.38), Inches(0.05), Inches(0.28),
                         size=Pt(9), color=DIM)

make_mixin_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — TRAINER INVENTORY
# ─────────────────────────────────────────────────────────────────────────────
def make_trainers_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Trainer Inventory",
                 "All 22 composed trainers — low-res, full-res, cascade-fullres stages")
    slide_number(sl, 10)

    def status_badge(sl, txt, x, y, color):
        add_rect(sl, x, y, Inches(0.88), Inches(0.26),
                 fill=RGBColor(*(int(c*0.2) for c in color)), line=color)
        add_text(sl, txt, x, y + Inches(0.03), Inches(0.88), Inches(0.22),
                 size=Pt(8.5), bold=True, color=color, align=PP_ALIGN.CENTER)

    # Low-res / Full-res column
    groups = [
        ("Full-Res (3d_fullres)", CYAN, [
            ("nnUNetTrainerDA5", "Baseline", GREEN, "Working"),
            ("DA5FiLMV3", "DA5 + FiLM bottleneck (K=8)", GREEN, "Working"),
            ("DA5FiLMTopo", "DA5 + FiLM + clDice", YELLOW, "Untested"),
            ("DA5TopoScheduled", "DA5 + clDice (scheduled)", PURPLE, "New"),
            ("DA5FiLMTopoCurriculum", "DA5 + FiLM + clDice + Curriculum", YELLOW, "Untested"),
        ]),
        ("Low-Res / Cascade Stage 1 (3d_lowres)", YELLOW, [
            ("DA5CascadeFiLM", "DA5 + FiLM", GREEN, "Working"),
            ("DA5CascadeTopo", "DA5 + clDice (fixed)", GREEN, "Working"),
            ("DA5CascadeFiLMTopo", "DA5 + FiLM + clDice", GREEN, "Working"),
            ("DA5CascadeTopoScheduled", "DA5 + clDice (scheduled)", PURPLE, "New"),
            ("DA5CascadeFiLMAdjacency", "DA5 + FiLM + adjacency", RED, "Pending"),
        ]),
        ("High-Res / Cascade Stage 2 (3d_cascade_fullres)", PURPLE, [
            ("DA5CascadeFullresBaseline", "DA5 only (name wrapper)", GREEN, "Working"),
            ("DA5CascadeFullresFiLM", "DA5 + FiLM", GREEN, "Working"),
            ("DA5CascadeFullresTopo", "DA5 + clDice  ← FIXED", PURPLE, "Fixed"),
            ("DA5CascadeFullresFiLMTopo", "DA5 + FiLM + clDice  ← FIXED", PURPLE, "Fixed"),
            ("DA5CascadeFullresTopoScheduled", "DA5 + clDice (scheduled)", PURPLE, "New"),
        ]),
    ]

    status_colors = {"Working": GREEN, "Untested": YELLOW, "New": PURPLE,
                     "Fixed": CYAN, "Pending": RED}

    cx = Inches(0.35)
    cy = Inches(1.15)
    col_w = Inches(4.15)

    for gi, (gtitle, gcolor, trainers) in enumerate(groups):
        gx = cx + gi * (col_w + Inches(0.2))
        add_rect(sl, gx, cy, col_w, Inches(0.38),
                 fill=RGBColor(0x14,0x18,0x2C), line=gcolor)
        add_text(sl, gtitle, gx + Inches(0.1), cy + Inches(0.05), col_w - Inches(0.2), Inches(0.3),
                 size=Pt(10), bold=True, color=gcolor)

        for ti, (name, desc, scolor, status) in enumerate(trainers):
            ty = cy + Inches(0.44) + ti * Inches(1.12)
            add_rect(sl, gx, ty, col_w, Inches(1.05),
                     fill=RGBColor(0x12, 0x16, 0x26) if ti % 2 == 0 else BG_DARK,
                     line=BORDER)
            add_text(sl, name, gx + Inches(0.1), ty + Inches(0.06),
                     col_w - Inches(1.1), Inches(0.32),
                     size=Pt(10.5), bold=True, color=CYAN, font_name="Courier New")
            sc = status_colors.get(status, MUTED)
            status_badge(sl, status, gx + col_w - Inches(0.95), ty + Inches(0.06), sc)
            add_text(sl, desc, gx + Inches(0.1), ty + Inches(0.44),
                     col_w - Inches(0.2), Inches(0.5),
                     size=Pt(10), color=MUTED)
            add_text(sl, "_100epochs variant available",
                     gx + Inches(0.1), ty + Inches(0.79), col_w - Inches(0.2), Inches(0.22),
                     size=Pt(8.5), color=DIM, italic=True)

make_trainers_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CASCADE ABLATION MATRIX
# ─────────────────────────────────────────────────────────────────────────────
def make_ablation_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Cascade Ablation Experiments",
                 "Systematic comparison of FiLM conditioning × topology loss at low-res stage")
    slide_number(sl, 11)

    exps = [
        ("Baseline", "nnUNetTrainerDA5_100epochs",             "DA5CascadeFullresBaseline", False, False, BG_DARK),
        ("A",        "DA5CascadeFiLM_100epochs",               "DA5CascadeFullresFiLM",     True,  False, RGBColor(0x14,0x18,0x2C)),
        ("B",        "DA5CascadeFiLMTopo_100epochs",           "DA5CascadeFullresFiLM",     True,  True,  RGBColor(0x14,0x18,0x2C)),
        ("D",        "DA5CascadeTopo_100epochs",               "DA5CascadeFullresTopo",     False, True,  RGBColor(0x14,0x18,0x2C)),
        ("B+",       "DA5CascadeFiLMTopo_100epochs",           "DA5CascadeFullresFiLMTopo", True,  True,  RGBColor(0x16,0x14,0x2C)),
        ("D-S",      "DA5CascadeTopoScheduled_100epochs",      "DA5CascadeFullresTopoScheduled", False, True, RGBColor(0x14,0x18,0x2C)),
    ]

    headers = ["Exp", "Low-res trainer", "High-res trainer", "FiLM", "Topo loss", "Mean Dice", "AO Dice", "PA Dice"]
    col_w   = [Inches(0.45), Inches(3.5), Inches(3.5), Inches(0.65), Inches(0.8), Inches(1.0), Inches(1.0), Inches(1.0)]
    col_x   = [Inches(0.35)]
    for w in col_w[:-1]: col_x.append(col_x[-1] + w)

    Y0 = Inches(1.15)
    ROW_H = Inches(0.6)

    # Header
    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_rect(sl, cx, Y0, cw, ROW_H - Inches(0.06),
                 fill=RGBColor(0x20,0x22,0x40), line=BORDER)
        add_text(sl, hdr, cx + Inches(0.05), Y0 + Inches(0.1), cw - Inches(0.1), ROW_H,
                 size=Pt(10), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    for ri, (label, lowres, highres, film, topo, bg) in enumerate(exps):
        ry = Y0 + ROW_H * (ri + 1)
        for ci, (val, cx, cw) in enumerate(zip(
                [label, lowres, highres,
                 "✓" if film else "—",
                 "✓" if topo else "—",
                 "—", "—", "—"],
                col_x, col_w)):
            fill = bg
            lc = BORDER
            if ci == 0:
                fill = RGBColor(0x20,0x22,0x40)
                lc = PURPLE if label in ("B+","D-S") else BORDER
            elif ci == 3:
                fill = RGBColor(0x14,0x22,0x2A) if film else bg
            elif ci == 4:
                fill = RGBColor(0x1A,0x14,0x30) if topo else bg
            elif ci >= 5:
                fill = RGBColor(0x14,0x1A,0x14)

            add_rect(sl, cx, ry, cw, ROW_H - Inches(0.04), fill=fill, line=lc)
            vc = CYAN if ci == 1 else (MUTED if ci == 2 else
                 (GREEN if val == "✓" else (RED if val == "—" and ci in (3,4) else
                 (YELLOW if ci >= 5 else WHITE))))
            add_text(sl, str(val), cx + Inches(0.04), ry + Inches(0.12),
                     cw - Inches(0.08), ROW_H - Inches(0.2),
                     size=Pt(9.5) if ci in (1,2) else Pt(11),
                     color=vc, align=PP_ALIGN.CENTER, bold=(ci == 0))

    # Legend
    legend = [("NEW/FIXED", PURPLE), ("FiLM", GREEN), ("Topo", PURPLE), ("Results TBD", YELLOW)]
    for i, (lbl, col) in enumerate(legend):
        lx = Inches(0.35) + i * Inches(2.0)
        add_rect(sl, lx, Inches(5.65), Inches(0.16), Inches(0.16), fill=col)
        add_text(sl, lbl, lx + Inches(0.22), Inches(5.64), Inches(1.7), Inches(0.26),
                 size=Pt(10), color=MUTED)

    add_text(sl, "B+ and D-S use the newly fixed/created high-res trainers that actually apply topology loss at both stages.",
             Inches(0.35), Inches(6.05), Inches(12.6), Inches(0.35),
             size=Pt(10.5), color=MUTED, italic=True)
    add_text(sl, "All 100-epoch variants. Use scripts/setup_cascade_predictions.py to create lowres prediction symlinks before training stage 2.",
             Inches(0.35), Inches(6.45), Inches(12.6), Inches(0.35),
             size=Pt(10.5), color=MUTED, italic=True)

make_ablation_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — NEW CONTRIBUTIONS (THIS SESSION)
# ─────────────────────────────────────────────────────────────────────────────
def make_contributions_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Contributions — This Implementation Session",
                 "Three changes to the codebase")
    slide_number(sl, 12)

    items = [
        (
            "1",
            "TopologyLossScheduledMixin  — NEW",
            "variants/mixins/topology_loss.py",
            PURPLE,
            [
                "Subclasses TopologyLossMixin — inherits all clDice logic",
                "Adds mixin_on_train_epoch_start() to update self.topo_weight each epoch",
                "Uses existing topo_weight_schedule() function (was previously unused)",
                "Schedule: warmup=10 epochs, plateau → decay from epoch 40",
                "Peak weight 0.5 (vs fixed 0.2) — more aggressive topology enforcement at plateau",
                "All schedule params overridable as class attributes in thin subclasses",
            ]
        ),
        (
            "2",
            "Fixed high-res cascade topology trainers",
            "composed/nnUNetTrainerDA5CascadeFullresVariants.py",
            CYAN,
            [
                "DA5CascadeFullresTopo was previously a name-only wrapper around plain DA5",
                "DA5CascadeFullresFiLMTopo was FiLM-only — topology was silently missing",
                "Both now properly use TopologyLossMixin at the high-res stage",
                "Enables end-to-end topology enforcement: low-res (global) + high-res (local)",
                "⚠  Checkpoints from previous runs with these names are incompatible",
            ]
        ),
        (
            "3",
            "New scheduled topology composed trainers  — NEW",
            "composed/nnUNetTrainerDA5TopoScheduled.py  (new file)",
            GREEN,
            [
                "DA5TopoScheduled(_100epochs) — full-res with scheduled clDice",
                "DA5CascadeTopoScheduled(_100epochs) — low-res cascade with scheduled clDice",
                "DA5CascadeFullresTopoScheduled(_100epochs) — high-res cascade with scheduled clDice",
                "All follow the same thin composed trainer pattern (~10 lines each)",
                "Pair DA5CascadeTopoScheduled (low-res) + DA5CascadeFullresTopoScheduled (high-res)",
            ]
        ),
    ]

    for i, (num, title, path, color, bullets) in enumerate(items):
        y = Inches(1.18) + i * Inches(1.98)
        add_rect(sl, Inches(0.35), y, Inches(12.6), Inches(1.9),
                 fill=RGBColor(0x12,0x16,0x26), line=color)
        add_rect(sl, Inches(0.35), y, Inches(12.6), Inches(0.04), fill=color)

        # Number badge
        add_rect(sl, Inches(0.35), y, Inches(0.5), Inches(0.5),
                 fill=color)
        add_text(sl, num, Inches(0.35), y + Inches(0.06), Inches(0.5), Inches(0.38),
                 size=Pt(18), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        add_text(sl, title, Inches(0.95), y + Inches(0.1), Inches(8), Inches(0.4),
                 size=Pt(13.5), bold=True, color=color)
        add_text(sl, path, Inches(0.95), y + Inches(0.5), Inches(10), Inches(0.3),
                 size=Pt(10.5), color=MUTED, italic=True, font_name="Courier New")

        for j, bullet in enumerate(bullets[:3]):
            add_text(sl, f"◆  {bullet}",
                     Inches(0.95), y + Inches(0.82) + j * Inches(0.32),
                     Inches(5.6), Inches(0.3),
                     size=Pt(10), color=WHITE)
        for j, bullet in enumerate(bullets[3:]):
            add_text(sl, f"◆  {bullet}",
                     Inches(6.85), y + Inches(0.82) + j * Inches(0.32),
                     Inches(5.9), Inches(0.3),
                     size=Pt(10), color=WHITE)

make_contributions_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — OPEN QUESTIONS
# ─────────────────────────────────────────────────────────────────────────────
def make_questions_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Open Questions & Assumptions",
                 "Items to validate before drawing conclusions from experiment results")
    slide_number(sl, 13)

    questions = [
        (YELLOW, "Q1", "Does FiLM actually learn?",
         "Monitor bottleneck_film.weight.abs().mean() during training.\nγ should diverge from 0 after first few epochs."),
        (YELLOW, "Q2", "clDice at high-res: local vs global?",
         "High-res patches cover only part of AO/PA.\nTruncated skeleton may give degraded clDice signal."),
        (ORANGE, "Q3", "Is topo_w_high = 0.5 too aggressive?",
         "Peak weight 0.5 gives topology ~33% of total gradient.\nMonitor training loss curves on first scheduled run."),
        (YELLOW, "Q4", "Memory cost of 10 skeleton iters in 3D?",
         "20 pooling ops per batch (2 classes × 10 iters).\nConsider num_iter=5 if GPU memory is constrained."),
        (RED, "Q5", "HLHS / PuA edge cases",
         "Absent/fused vessels → near-empty GT skeleton.\nTopology loss may mis-penalise correct empty predictions."),
        (YELLOW, "Q6", "Cascade symlink script coverage",
         "setup_cascade_predictions.py must handle all new\ntrainer name combos (Topo, FiLMTopo, TopoScheduled)."),
        (YELLOW, "Q7", "disease_map.json at cascade-fullres inference",
         "FiLM at high-res stage needs disease vec at inference.\nVerify predict_disease_conditioned.py handles cascade."),
        (ORANGE, "Q8", "DA5 augmentation for cardiac CT",
         "Aggressive deformation may be unrealistic for rigid\ncardiac structures. Compare DA5 vs default augmentation."),
    ]

    for i, (color, qid, title, desc) in enumerate(questions):
        col = i // 4
        row = i % 4
        qx = Inches(0.35) + col * Inches(6.5)
        qy = Inches(1.18) + row * Inches(1.48)
        add_rect(sl, qx, qy, Inches(6.25), Inches(1.38),
                 fill=RGBColor(0x14,0x18,0x2C), line=color)
        add_rect(sl, qx, qy, Inches(6.25), Inches(0.04), fill=color)
        add_rect(sl, qx, qy, Inches(0.55), Inches(0.42), fill=color)
        add_text(sl, qid, qx, qy + Inches(0.04), Inches(0.55), Inches(0.34),
                 size=Pt(11), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(sl, title, qx + Inches(0.62), qy + Inches(0.08), Inches(5.5), Inches(0.34),
                 size=Pt(12), bold=True, color=color)
        add_text(sl, desc, qx + Inches(0.12), qy + Inches(0.5), Inches(6.0), Inches(0.8),
                 size=Pt(10.5), color=MUTED)

make_questions_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — FUTURE DIRECTIONS
# ─────────────────────────────────────────────────────────────────────────────
def make_future_slide():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_DARK)
    slide_header(sl, "Future Directions",
                 "Extensions to explore once training protocol is validated")
    slide_number(sl, 14)

    futures = [
        (PURPLE, "Per-Disease Topology Exclusion",
         "TopologyLossConditionalMixin: disable clDice for HLHS (flag 0) and PuA (flag 5) "
         "cases where AO/PA are absent. Reads disease vector from batch to gate the loss."),
        (CYAN, "Adjacency Loss for VSD/AVSD",
         "DA5CascadeFiLMAdjacency is a placeholder. An adjacency-consistency loss penalising "
         "disconnected LV–RV in VSD cases would complement FiLM conditioning."),
        (GREEN, "clDice as Primary Loss",
         "Replace Dice+CE with α·Dice+(1-α)·clDice as the main loss (not extra term). "
         "Original paper shows stronger topology preservation. New ClDiceMainLossMixin."),
        (YELLOW, "Multi-Stage Encoder FiLM",
         "Apply FiLM at 2 lowest-resolution encoder stages instead of bottleneck-only. "
         "Larger receptive field, more spatial context, lower N so (1+γ)^N instability avoided."),
        (ORANGE, "Cross-Attention Conditioning",
         "Replace FiLM scale/shift with cross-attention: disease vector attends to spatial "
         "bottleneck features. Selective amplification vs. uniform channel rescaling."),
        (RED, "Post-Processing: Disease-Aware Largest-CC",
         "Rule-based cleanup: keep N largest connected components per label based on disease "
         "flags. Fast, interpretable. TGA/DORV: enforce single AO/PA connected component."),
    ]

    for i, (color, title, desc) in enumerate(futures):
        col = i % 3
        row = i // 3
        fx = Inches(0.35) + col * Inches(4.35)
        fy = Inches(1.18) + row * Inches(2.6)
        add_rect(sl, fx, fy, Inches(4.15), Inches(2.45),
                 fill=RGBColor(0x12,0x16,0x26), line=color)
        add_rect(sl, fx, fy, Inches(4.15), Inches(0.05), fill=color)
        add_text(sl, title, fx + Inches(0.15), fy + Inches(0.1), Inches(3.85), Inches(0.42),
                 size=Pt(13), bold=True, color=color)
        add_text(sl, desc, fx + Inches(0.15), fy + Inches(0.58), Inches(3.85), Inches(1.75),
                 size=Pt(11), color=MUTED)

make_future_slide()

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "..", "docs", "CHD_TopologyLoss_Presentation.pptx")
out = os.path.normpath(out)
prs.save(out)
print(f"Saved → {out}")
